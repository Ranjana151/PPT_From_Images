# PPt from images after apply logo on each image  by maintain aspect ratio of original images
import os
from wand.image import Image
from pptx import Presentation
from pptx.util import Inches

# Original Image resizing by maintaining aspect ratio

def resize_original_image(image1, new_width):
    width, height = image1.size
    aspect_ratio = width/height
    new_height = int(aspect_ratio*new_width)
    image1.resize(new_width, new_height)
    return image1

# Logo Image resizing
def resize_original_logo(image2, new_width):
    width, height = image2.size
    aspect_ratio = width/height
    new_height = 15
    image2.resize(new_width, new_height)
    return image2


# list all  images from a given folder
image_directory = "Resources"
image_list = os.listdir(image_directory)
image_list = image_list[0:5]



prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

with Image(filename="Resources/nike_black.png") as img1:
    img1 = resize_original_logo(img1, 75)
    for image in image_list:
        Resource_file = "Resources/"+image
        water_mark_image_path = "output/"+image
        with Image(filename=Resource_file) as img2:
            img2 = resize_original_image(img2, 200)

            # Apply watermark process on original image one by one
            img2.composite_channel('all_channels', img1, 'dissolve', 4, 4)

            # Adding watermark images on ppt slide one by one
            img2.save(filename=water_mark_image_path)
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = 'Sample Title 1'
            tf = body_shape.text_frame
            tf.text = 'Sample Subtitle 1'
            left = Inches(1)
            top = Inches(2.8)
            pic = slide.shapes.add_picture(water_mark_image_path, left, top)
prs.save('Resources/watermark_images.pptx')











